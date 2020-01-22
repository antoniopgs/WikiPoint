from pptx import Presentation
import wikipedia
import re

def search():
    search_term = input("Search: ")
    # Fix Search Term in case of errors:
    fixed_search_term = wikipedia.suggest(search_term)
    # If Search Term needs fixing:
    if fixed_search_term:
        return wikipedia.search(fixed_search_term)
    # If Search Term doesn't need fixing:
    else:
        return wikipedia.search(search_term)

def select_topic(search_term):
    while True:
        print("---------- DETECTED TOPICS ---------")
        # Iterate over Search Result's Length and Contents at the same time:
        for i, topic in zip(range(len(search_term)), search_term):
            print(f"{i+1} - {topic}")
        topic_number = input("\nInsert Topic Number: ")
        try: 
            if 1 <= int(topic_number) <= len(search_term):
                return int(topic_number)
            else:
                print("Invalid Topic Number.\n")
        except TypeError:
            print("Invalid Topic Number.\n")

# Search and get info:
search_result = search()
topic_result = select_topic(search_result)
page = wikipedia.page(search_result[topic_result-1])
pre_info = page.content.split("\n== See also ==\n")
info = pre_info[0].strip()

# Write info to .txt file:
with open("text.txt", "w", encoding="utf-8") as file:
    file.write(info)

# Find all tags and add page title as first tag:
tags = [tag.strip().title() for tag in re.findall("\n==.*==\n", info)]
tags.insert(0, page.title)

# Split by all tags. Get tag's texts:
tags_text = [section.strip() for section in re.split("\n==.*==\n", info)]

# Create List with all Tags and Text:
sections = list(zip(tags, tags_text))

# Create PowerPoint:
ppt = Presentation()
# Slide Templates:
just_title = ppt.slide_layouts[0]
title_content = ppt.slide_layouts[1]
section_header = ppt.slide_layouts[2]
two_content = ppt.slide_layouts[3]
comparison = ppt.slide_layouts[4]
title_only = ppt.slide_layouts[5]
blank = ppt.slide_layouts[6]
content_caption = ppt.slide_layouts[7]
picture_caption = ppt.slide_layouts[8]
# Generator:
for i in range(len(sections)):
    if i == 0:
        # Cover Slide:
        cover_slide = ppt.slides.add_slide(just_title)
        cover_slide_title = cover_slide.shapes.title
        cover_slide_title.text = sections[0][0]
        # Index Slide:
        index_slide = ppt.slides.add_slide(title_content)
        index_slide_title = index_slide.shapes.title
        index_slide_title.text = "Index"
        # Index Generator:
        index = index_slide.shapes.placeholders[1]
        tf = index.text_frame
        #tf.text = "Find the bullet slide layout"
        for section in sections[1:]:
            p = tf.add_paragraph()
            if re.match("={2} .* ={2}", section[0]):
                p.level = 1
            elif re.match("={3} .* ={3}", section[0]):
                p.level = 2
            elif re.match("={4} .* ={4}", section[0]):
                p.level = 3
            elif re.match("={5} .* ={5}", section[0]):
                p.level = 4
            elif re.match("={6} .* ={6}", section[0]):
                p.level = 5
            p.text = section[0].replace("=", "").strip()
    else:
        info_slide = ppt.slides.add_slide(section_header)
        title = info_slide.shapes.title
        title.text = sections[i][0]
ppt.save('presentation.pptx')